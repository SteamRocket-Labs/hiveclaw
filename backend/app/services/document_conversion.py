from __future__ import annotations

import hashlib
import json
import mimetypes
import re
from dataclasses import dataclass
from datetime import UTC, datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Literal

from loguru import logger

ConversionMode = Literal["auto", "fast", "ocr", "layout", "vision"]
ReturnFormat = Literal["preview", "markdown", "metadata", "pages"]

_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".js",
    ".ts",
    ".py",
    ".html",
    ".htm",
    ".css",
    ".sh",
    ".log",
    ".env",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".sql",
}


@dataclass(frozen=True)
class DocumentConversionRequest:
    source_path: Path
    workspace_root: Path
    source_uri: str | None
    tenant_id: object | None
    agent_id: object | None
    user_id: object | None
    mode: ConversionMode = "auto"
    max_pages: int | None = None
    max_output_chars: int | None = None
    force_refresh: bool = False


@dataclass(frozen=True)
class DocumentConversionResult:
    markdown: str
    plain_text: str
    source_path: str
    source_uri: str | None
    source_sha256: str
    source_mime_type: str
    engine: str
    used_ocr: bool
    used_vision: bool
    page_count: int | None
    artifact_markdown_path: str
    artifact_metadata_path: str
    warnings: tuple[str, ...]


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._skip_depth = 0
        self._parts: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in {"script", "style", "noscript"}:
            self._skip_depth += 1
        if tag in {"p", "div", "section", "article", "br", "li", "tr", "h1", "h2", "h3", "h4"}:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"} and self._skip_depth:
            self._skip_depth -= 1
        if tag in {"p", "div", "section", "article", "li", "tr", "h1", "h2", "h3", "h4"}:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip_depth and data.strip():
            self._parts.append(data.strip())

    def get_text(self) -> str:
        text = " ".join(part.strip() for part in self._parts if part.strip())
        return re.sub(r"\s*\n\s*", "\n", text).strip()


class DocumentConversionService:
    def __init__(self, *, max_file_bytes: int = 50 * 1024 * 1024) -> None:
        self.max_file_bytes = max_file_bytes

    def convert(self, request: DocumentConversionRequest) -> DocumentConversionResult:
        workspace_root = request.workspace_root.resolve()
        source_path = request.source_path.resolve()
        self._ensure_inside_workspace(source_path, workspace_root)
        data = source_path.read_bytes()
        return self._convert_data(
            data=data,
            source_path=source_path,
            workspace_root=workspace_root,
            source_uri=request.source_uri,
            source_mime_type=mimetypes.guess_type(source_path.name)[0] or "application/octet-stream",
            tenant_id=request.tenant_id,
            agent_id=request.agent_id,
            user_id=request.user_id,
            mode=request.mode,
            max_pages=request.max_pages,
            force_refresh=request.force_refresh,
        )

    def convert_bytes(
        self,
        *,
        data: bytes,
        filename: str,
        workspace_root: Path,
        source_uri: str | None = None,
        source_mime_type: str | None = None,
        tenant_id: object | None = None,
        agent_id: object | None = None,
        user_id: object | None = None,
        mode: ConversionMode = "auto",
        max_pages: int | None = None,
        force_refresh: bool = False,
    ) -> DocumentConversionResult:
        workspace_root = workspace_root.resolve()
        workspace_root.mkdir(parents=True, exist_ok=True)
        sha = _sha256_bytes(data)
        safe_name = _safe_filename(filename)
        source_dir = self._hive_base_dir(workspace_root, None) / "document_sources" / sha
        source_dir.mkdir(parents=True, exist_ok=True)
        source_path = source_dir / safe_name
        if force_refresh or not source_path.exists() or source_path.read_bytes() != data:
            source_path.write_bytes(data)
        return self._convert_data(
            data=data,
            source_path=source_path,
            workspace_root=workspace_root,
            source_uri=source_uri,
            source_mime_type=source_mime_type or mimetypes.guess_type(safe_name)[0] or "application/octet-stream",
            tenant_id=tenant_id,
            agent_id=agent_id,
            user_id=user_id,
            mode=mode,
            max_pages=max_pages,
            force_refresh=force_refresh,
        )

    def _convert_data(
        self,
        *,
        data: bytes,
        source_path: Path,
        workspace_root: Path,
        source_uri: str | None,
        source_mime_type: str,
        tenant_id: object | None,
        agent_id: object | None,
        user_id: object | None,
        mode: ConversionMode,
        max_pages: int | None,
        force_refresh: bool,
    ) -> DocumentConversionResult:
        if len(data) > self.max_file_bytes:
            raise ValueError(f"File too large for document conversion: {len(data)} bytes")

        source_sha256 = _sha256_bytes(data)
        hive_base = self._hive_base_dir(workspace_root, source_path)
        artifact_dir = hive_base / "document_conversions" / source_sha256
        markdown_path = artifact_dir / "content.md"
        metadata_path = artifact_dir / "metadata.json"
        markdown_rel = _relative_path(markdown_path, workspace_root)
        metadata_rel = _relative_path(metadata_path, workspace_root)

        if not force_refresh and markdown_path.exists() and metadata_path.exists():
            metadata = _read_json(metadata_path)
            if metadata.get("source_uri") == source_uri:
                markdown = markdown_path.read_text(encoding="utf-8", errors="replace")
                return DocumentConversionResult(
                    markdown=markdown,
                    plain_text=_plain_text(markdown),
                    source_path=str(source_path),
                    source_uri=source_uri,
                    source_sha256=source_sha256,
                    source_mime_type=metadata.get("source_mime_type") or source_mime_type,
                    engine=metadata.get("engine") or "unknown",
                    used_ocr=bool(metadata.get("used_ocr")),
                    used_vision=bool(metadata.get("used_vision")),
                    page_count=metadata.get("page_count"),
                    artifact_markdown_path=markdown_rel,
                    artifact_metadata_path=metadata_rel,
                    warnings=tuple(metadata.get("warnings") or ()),
                )

        markdown, engine, warnings = self._convert_source(source_path, data, source_mime_type, source_uri)
        markdown = markdown.strip()
        if not markdown:
            warnings.append("empty_conversion_output")
            markdown = ""

        artifact_dir.mkdir(parents=True, exist_ok=True)
        markdown_path.write_text(markdown, encoding="utf-8")
        metadata = {
            "source_path": str(source_path),
            "source_uri": source_uri,
            "source_sha256": source_sha256,
            "source_mime_type": source_mime_type,
            "engine": engine,
            "used_ocr": False,
            "used_vision": False,
            "page_count": None,
            "artifact_markdown_path": markdown_rel,
            "artifact_metadata_path": metadata_rel,
            "warnings": warnings,
            "tenant_id": str(tenant_id) if tenant_id is not None else None,
            "agent_id": str(agent_id) if agent_id is not None else None,
            "user_id": str(user_id) if user_id is not None else None,
            "mode": mode,
            "max_pages": max_pages,
            "created_at": datetime.now(UTC).isoformat(),
        }
        metadata_path.write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")

        return DocumentConversionResult(
            markdown=markdown,
            plain_text=_plain_text(markdown),
            source_path=str(source_path),
            source_uri=source_uri,
            source_sha256=source_sha256,
            source_mime_type=source_mime_type,
            engine=engine,
            used_ocr=False,
            used_vision=False,
            page_count=None,
            artifact_markdown_path=markdown_rel,
            artifact_metadata_path=metadata_rel,
            warnings=tuple(warnings),
        )

    def _convert_source(
        self, source_path: Path, data: bytes, source_mime_type: str, source_uri: str | None
    ) -> tuple[str, str, list[str]]:
        warnings: list[str] = []
        is_pdf = _is_pdf_source(source_path, data, source_mime_type)
        try:
            markdown = self._convert_with_markitdown(source_path)
            stripped_markdown = markdown.strip()
            if is_pdf and _looks_like_raw_pdf_markdown(stripped_markdown):
                warnings.append("local_markitdown_unreadable_pdf")
            elif stripped_markdown:
                return markdown, "local_markitdown", warnings
            warnings.append("local_markitdown_empty_output")
        except ModuleNotFoundError:
            warnings.append("local_markitdown_unavailable")
        except Exception as exc:
            logger.warning("[DocumentConversion] MarkItDown conversion failed for %s: %s", source_path, exc)
            warnings.append(f"local_markitdown_failed:{type(exc).__name__}")

        legacy = self._convert_with_legacy_extractors(source_path, data, source_mime_type, source_uri)
        return legacy, "legacy_python_extractors", warnings

    def _convert_with_markitdown(self, source_path: Path) -> str:
        from markitdown import MarkItDown

        try:
            converter = MarkItDown(enable_plugins=False)
        except TypeError:
            converter = MarkItDown()

        if hasattr(converter, "convert_local"):
            result = converter.convert_local(str(source_path))
        else:
            result = converter.convert(str(source_path))
        return str(getattr(result, "text_content", "") or "")

    def _convert_with_legacy_extractors(
        self, source_path: Path, data: bytes, source_mime_type: str, source_uri: str | None
    ) -> str:
        ext = source_path.suffix.lower()
        if ext in _TEXT_EXTENSIONS or source_mime_type.startswith("text/"):
            text = _decode_text(data)
            if ext in {".html", ".htm"} or "html" in source_mime_type:
                return self._extract_html_text(text, source_uri=source_uri)
            return text
        if ext == ".pdf" or source_mime_type == "application/pdf" or data[:5].startswith(b"%PDF"):
            return _extract_pdf(data)
        if ext == ".docx":
            return _extract_docx(data)
        if ext == ".xlsx":
            return _extract_xlsx(data)
        if ext == ".pptx":
            return _extract_pptx(data)
        return ""

    def _extract_html_text(self, markup: str, *, source_uri: str | None = None) -> str:
        try:
            import trafilatura

            extracted = trafilatura.extract(
                markup.strip(),
                url=source_uri,
                include_comments=False,
                include_tables=True,
                favor_precision=True,
            )
            if extracted and extracted.strip():
                return extracted.strip()
        except Exception as exc:
            logger.debug("[DocumentConversion] Trafilatura unavailable/failed: %s", exc)

        parser = _HTMLTextExtractor()
        parser.feed(markup)
        return parser.get_text()

    def _hive_base_dir(self, workspace_root: Path, source_path: Path | None) -> Path:
        workspace_dir = workspace_root / "workspace"
        if source_path is not None:
            try:
                source_path.resolve().relative_to(workspace_dir.resolve())
                return workspace_dir / ".hive"
            except ValueError:
                pass
        return workspace_root / ".hive"

    @staticmethod
    def _ensure_inside_workspace(source_path: Path, workspace_root: Path) -> None:
        try:
            source_path.resolve().relative_to(workspace_root.resolve())
        except ValueError as exc:
            raise ValueError("Document conversion source is outside workspace") from exc


def render_conversion_preview(result: DocumentConversionResult, *, max_chars: int = 8000) -> str:
    preview = result.markdown.strip()
    if len(preview) > max_chars:
        preview = preview[:max_chars] + f"\n\n...[truncated, {len(result.markdown)} chars total]"
    warnings = ""
    if result.warnings:
        warnings = "\nWarnings: " + ", ".join(result.warnings)
    return (
        f"Converted with {result.engine}.\n"
        f"Full Markdown: {result.artifact_markdown_path}\n"
        f"Metadata: {result.artifact_metadata_path}"
        f"{warnings}\n\n"
        "Preview:\n"
        f"{preview}"
    )


def _safe_filename(filename: str) -> str:
    safe = Path(filename.replace("\\", "/")).name.strip()
    if not safe or safe in {".", ".."}:
        return "source.bin"
    return re.sub(r"[^A-Za-z0-9._-]+", "_", safe)


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _relative_path(path: Path, root: Path) -> str:
    try:
        return path.resolve().relative_to(root.resolve()).as_posix()
    except ValueError:
        return str(path.resolve())


def _read_json(path: Path) -> dict:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _is_pdf_source(source_path: Path, data: bytes, source_mime_type: str) -> bool:
    return (
        source_path.suffix.lower() == ".pdf"
        or "application/pdf" in source_mime_type.lower()
        or data[:5].startswith(b"%PDF")
    )


def _looks_like_raw_pdf_markdown(markdown: str) -> bool:
    return markdown.lstrip().startswith("%PDF")


def _plain_text(markdown: str) -> str:
    text = re.sub(r"```.*?```", "", markdown, flags=re.DOTALL)
    text = re.sub(r"[*_`>#\[\]()]|!\[[^\]]*\]\([^)]*\)", "", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _extract_pdf(data: bytes) -> str:
    import io

    try:
        import pdfplumber

        parts: list[str] = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for index, page in enumerate(pdf.pages[:50]):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(f"--- Page {index + 1} ---\n{page_text.strip()}")
                for table in page.extract_tables() or []:
                    rows = []
                    for row in table:
                        cells = [str(cell or "").strip() for cell in row]
                        if any(cells):
                            rows.append(" | ".join(cells))
                    if rows:
                        parts.append("Table:\n" + "\n".join(rows))
        if parts:
            return "\n\n".join(parts)
    except ModuleNotFoundError:
        pass
    except Exception as exc:
        logger.debug("[DocumentConversion] pdfplumber extraction failed: %s", exc)

    try:
        import io

        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        parts = []
        for index, page in enumerate(reader.pages[:50]):
            page_text = page.extract_text() or ""
            if page_text.strip():
                parts.append(f"--- Page {index + 1} ---\n{page_text.strip()}")
        return "\n\n".join(parts)
    except Exception as exc:
        logger.debug("[DocumentConversion] PDF legacy extraction failed: %s", exc)
        return ""


def _extract_docx(data: bytes) -> str:
    import io

    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append("Table:\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import io

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for sheet_name in wb.sheetnames[:10]:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(max_row=200, values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        return "\n\n".join(parts)
    finally:
        wb.close()


def _extract_pptx(data: bytes) -> str:
    import io

    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides[:50]):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"--- Slide {index + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(parts)
