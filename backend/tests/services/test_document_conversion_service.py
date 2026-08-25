from __future__ import annotations

import json
import sys
from io import BytesIO
from types import SimpleNamespace

import pytest


def _install_fake_markitdown(monkeypatch, calls: list[dict], output: str = "# Converted\n\nBody") -> None:
    class _FakeMarkItDown:
        def __init__(self, *args, **kwargs):
            calls.append({"kind": "init", "args": args, "kwargs": kwargs})

        def convert_local(self, source: str):
            calls.append({"kind": "convert_local", "source": source})
            return SimpleNamespace(text_content=output)

    monkeypatch.setitem(sys.modules, "markitdown", SimpleNamespace(MarkItDown=_FakeMarkItDown))


async def test_run_killable_in_process_timeout_physically_kills_worker(tmp_path):
    """A timed-out child process is physically killed mid-execution.

    The pool is warmed first so first-process startup overhead cannot eat the
    timeout; the started marker proves the worker really began executing, and
    the absent completed marker proves it was killed rather than abandoned.
    """
    import anyio

    from app.services.document_conversion import run_killable_in_process

    assert await run_killable_in_process(len, (1, 2, 3), timeout_seconds=30) == 3

    started = tmp_path / "worker_started.marker"
    completed = tmp_path / "worker_completed.marker"
    code = (
        "import time\n"
        "from pathlib import Path\n"
        f"Path({str(started)!r}).write_text('started', encoding='utf-8')\n"
        "time.sleep(2)\n"
        f"Path({str(completed)!r}).write_text('completed', encoding='utf-8')\n"
    )
    with pytest.raises(TimeoutError):
        await run_killable_in_process(exec, code, timeout_seconds=0.5)
    assert started.exists()
    await anyio.sleep(2.5)
    assert not completed.exists()


def test_document_conversion_uses_markitdown_local_file_and_writes_artifacts(monkeypatch, tmp_path):
    from app.services.document_conversion import DocumentConversionRequest, DocumentConversionService

    calls: list[dict] = []
    _install_fake_markitdown(monkeypatch, calls)

    workspace_root = tmp_path / "agent"
    source = workspace_root / "workspace" / "uploads" / "report.html"
    source.parent.mkdir(parents=True)
    source.write_text("<h1>Report</h1>", encoding="utf-8")

    result = DocumentConversionService().convert(
        DocumentConversionRequest(
            source_path=source,
            workspace_root=workspace_root,
            source_uri="https://example.com/report",
            tenant_id=None,
            agent_id=None,
            user_id=None,
            mode="auto",
        )
    )

    assert result.markdown == "# Converted\n\nBody"
    assert result.engine == "local_markitdown"
    assert result.source_uri == "https://example.com/report"
    assert result.artifact_markdown_path.startswith("workspace/.hive/document_conversions/")

    artifact_path = workspace_root / result.artifact_markdown_path
    metadata_path = workspace_root / result.artifact_metadata_path
    assert artifact_path.read_text(encoding="utf-8") == "# Converted\n\nBody"

    metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
    assert metadata["engine"] == "local_markitdown"
    assert metadata["source_uri"] == "https://example.com/report"
    assert metadata["artifact_markdown_path"] == result.artifact_markdown_path

    convert_calls = [call for call in calls if call["kind"] == "convert_local"]
    assert convert_calls == [{"kind": "convert_local", "source": str(source)}]


def test_document_conversion_reuses_cached_artifact_without_reinvoking_engine(monkeypatch, tmp_path):
    from app.services.document_conversion import DocumentConversionRequest, DocumentConversionService

    calls: list[dict] = []
    _install_fake_markitdown(monkeypatch, calls, output="First conversion")

    workspace_root = tmp_path / "agent"
    source = workspace_root / "workspace" / "uploads" / "cached.txt"
    source.parent.mkdir(parents=True)
    source.write_text("cache me", encoding="utf-8")
    request = DocumentConversionRequest(
        source_path=source,
        workspace_root=workspace_root,
        source_uri=None,
        tenant_id=None,
        agent_id=None,
        user_id=None,
        mode="auto",
    )

    first = DocumentConversionService().convert(request)
    second = DocumentConversionService().convert(request)

    assert first.artifact_markdown_path == second.artifact_markdown_path
    assert second.markdown == "First conversion"
    assert [call["kind"] for call in calls].count("convert_local") == 1


def test_document_conversion_rejects_symlink_escape(tmp_path):
    from app.services.document_conversion import DocumentConversionRequest, DocumentConversionService

    workspace_root = tmp_path / "agent"
    workspace = workspace_root / "workspace"
    workspace.mkdir(parents=True)
    outside = tmp_path / "outside.txt"
    outside.write_text("secret", encoding="utf-8")
    escaped = workspace / "escaped.txt"
    escaped.symlink_to(outside)

    with pytest.raises(ValueError, match="outside workspace"):
        DocumentConversionService().convert(
            DocumentConversionRequest(
                source_path=escaped,
                workspace_root=workspace_root,
                source_uri=None,
                tenant_id=None,
                agent_id=None,
                user_id=None,
                mode="auto",
            )
        )


def test_document_conversion_falls_back_to_legacy_text_when_markitdown_missing(monkeypatch, tmp_path):
    from app.services.document_conversion import DocumentConversionRequest, DocumentConversionService

    monkeypatch.setitem(sys.modules, "markitdown", None)
    workspace_root = tmp_path / "agent"
    source = workspace_root / "workspace" / "uploads" / "note.md"
    source.parent.mkdir(parents=True)
    source.write_text("# Native Markdown\n\nAlready useful.", encoding="utf-8")

    result = DocumentConversionService().convert(
        DocumentConversionRequest(
            source_path=source,
            workspace_root=workspace_root,
            source_uri=None,
            tenant_id=None,
            agent_id=None,
            user_id=None,
            mode="auto",
        )
    )

    assert result.engine == "legacy_python_extractors"
    assert "# Native Markdown" in result.markdown
    assert "local_markitdown_unavailable" in result.warnings


def test_document_conversion_rejects_raw_pdf_header_output_from_markitdown(monkeypatch, tmp_path):
    from app.services.document_conversion import DocumentConversionRequest, DocumentConversionService

    calls: list[dict] = []
    _install_fake_markitdown(monkeypatch, calls, output="%PDF-1.4\nthis is not extracted page text")

    workspace_root = tmp_path / "agent"
    source = workspace_root / "workspace" / "uploads" / "broken.pdf"
    source.parent.mkdir(parents=True)
    source.write_bytes(b"%PDF-1.4\nthis is not a parseable pdf body at all")

    result = DocumentConversionService().convert(
        DocumentConversionRequest(
            source_path=source,
            workspace_root=workspace_root,
            source_uri="https://issuer.example/broken.pdf",
            tenant_id=None,
            agent_id=None,
            user_id=None,
            mode="auto",
        )
    )

    assert result.markdown == ""
    assert result.engine == "legacy_python_extractors"
    assert "local_markitdown_unreadable_pdf" in result.warnings


def test_save_extracted_text_compatibility_wrapper_uses_document_conversion(monkeypatch, tmp_path):
    from app.services.text_extractor import save_extracted_text

    calls: list[dict] = []
    _install_fake_markitdown(monkeypatch, calls, output="# Compatibility\n\nConverted")

    workspace_root = tmp_path / "agent"
    source = workspace_root / "workspace" / "uploads" / "report.pdf"
    source.parent.mkdir(parents=True)
    source.write_bytes(b"%PDF fake")

    txt_path = save_extracted_text(source, source.read_bytes(), source.name)

    assert txt_path == source.parent / "report.txt"
    assert txt_path.read_text(encoding="utf-8") == "Compatibility\n\nConverted"
    assert (workspace_root / "workspace" / ".hive" / "document_conversions").exists()
    assert [call["kind"] for call in calls].count("convert_local") == 1


def test_legacy_pdf_extractor_preserves_pages_after_fifty(monkeypatch) -> None:
    from app.services.document_conversion import _extract_pdf

    class _Page:
        def __init__(self, text: str):
            self._text = text

        def extract_text(self):
            return self._text

        def extract_tables(self):
            return []

    class _Pdf:
        pages = [_Page(f"page-{index}") for index in range(51)]

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return None

    monkeypatch.setitem(sys.modules, "pdfplumber", SimpleNamespace(open=lambda _stream: _Pdf()))

    result = _extract_pdf(b"%PDF-fake")

    assert "page-50" in result


def test_legacy_xlsx_extractor_preserves_all_sheets_and_rows() -> None:
    from openpyxl import Workbook

    from app.services.document_conversion import _extract_xlsx

    workbook = Workbook()
    workbook.remove(workbook.active)
    for index in range(11):
        sheet = workbook.create_sheet(f"sheet-{index}")
        sheet.append([f"sheet-value-{index}"])
    for row_index in range(2, 202):
        workbook["sheet-10"].append([f"row-{row_index}"])
    stream = BytesIO()
    workbook.save(stream)

    result = _extract_xlsx(stream.getvalue())

    assert "## Sheet: sheet-10" in result
    assert "row-201" in result


def test_legacy_pptx_extractor_preserves_slides_after_fifty() -> None:
    from pptx import Presentation

    from app.services.document_conversion import _extract_pptx

    presentation = Presentation()
    for index in range(51):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"slide-{index}"
    stream = BytesIO()
    presentation.save(stream)

    result = _extract_pptx(stream.getvalue())

    assert "slide-50" in result
