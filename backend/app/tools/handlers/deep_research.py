from __future__ import annotations

import asyncio
import html
import json
import logging
import shutil
import uuid
from collections.abc import AsyncIterator
from io import BytesIO
from pathlib import Path
from typing import Any

from app.services.deep_research.artifact_composer import compose_deep_research_artifact
from app.services.deep_research.plan_mode import (
    build_deep_research_plan_fill,
    build_deep_research_plan_preview,
    deep_research_plan_signature,
    normalize_deep_research_output_format,
)
from app.services.deep_research.schemas import ResearchRequest, to_jsonable
from app.services.deep_research.workflow_definition import (
    start_deep_research_workflow_run,
)
from app.services.long_task_runtime import record_long_task_progress
from app.services.office_document_service import OfficeDocumentService
from app.services.plan_mode_service import get_plan_mode_service
from app.services.runtime_task_service import (
    get_runtime_task_record,
    update_runtime_task_record,
)
from app.tools.decorator import ToolMeta, tool
from app.tools.runtime import ToolExecutionRequest


logger = logging.getLogger(__name__)


def _schema(properties: dict[str, Any], required: list[str] | None = None) -> dict[str, Any]:
    return {"type": "object", "properties": properties, "required": required or []}


_REQUEST_PROPERTIES: dict[str, Any] = {
    "question": {"type": "string", "description": "The research question to answer."},
    "mode": {
        "type": "string",
        "enum": ["topic_deep_dive", "industry_research", "source_ledger_audit"],
        "description": "Research mode.",
    },
    "scope": {"type": "string", "description": "Optional geography, company, product, or excluded scope."},
    "depth": {"type": "string", "enum": ["quick", "standard", "full", "flagship"]},
    "source_policy": {"type": "string", "enum": ["primary_only", "primary_preferred", "mixed"]},
    "time_window": {"type": "string"},
    "max_rounds": {"type": "integer", "minimum": 1, "maximum": 8},
    "max_sources": {"type": "integer", "minimum": 1, "maximum": 50},
    "concurrency": {"type": "integer", "minimum": 1, "maximum": 12},
    "token_budget": {"type": "integer", "minimum": 1},
    "deadline_seconds": {"type": "integer", "minimum": 10},
    "output_format": {"type": "string", "enum": ["markdown", "json", "html", "docx", "xlsx", "pptx"]},
    "output_language": {
        "type": "string",
        "description": "Force the report's output language (e.g. 'zh', 'en', '中文'). Defaults to the question's language.",
    },
    "plan_confirmed": {
        "type": "boolean",
        "description": (
            "Set true ONLY after the user has reviewed and approved the research plan. Without it, the tool "
            "returns a plan + clarifying questions to confirm and does NOT execute the research."
        ),
    },
    "worker_topics": {
        "type": "array",
        "items": {"type": "string"},
        "description": "User-approved worker topics from the plan stage. Pass them back together with plan_confirmed=true.",
    },
}


# P3: best-effort in-process guard so the same question is not run twice concurrently
# (e.g. a stray deep_research_run + deep_research_start for the same ask).
_INFLIGHT_DEEP_RESEARCH: dict[tuple[str, str], str] = {}


def _clarifying_questions(research_request: ResearchRequest) -> list[str]:
    """Domain-general intent-alignment questions surfaced before any run."""
    return [
        "Scope & focus: which specific entities, products, regions, or sub-topics are in-scope, and what is explicitly out-of-scope?",
        f"Time window: currently '{research_request.time_window or 'unspecified'}' — confirm the period the research should cover.",
        f"Depth: currently '{research_request.depth}' — is that the right size, or do you want quick / full?",
        "Decision: what decision or deliverable should this research support? That sets the evidence bar.",
        "Output language: confirm the language the final report should be written in.",
    ]


async def _plan_preview(research_request: ResearchRequest) -> dict[str, Any]:
    """Build a deterministic plan + worker-topic preview (no LLM, no fan-out) for the plan gate."""
    return await build_deep_research_plan_preview(research_request)


async def _materialize_plan_card_payload(
    tool_request: ToolExecutionRequest,
    research_request: ResearchRequest,
    preview: dict[str, Any],
) -> dict[str, Any]:
    fill = build_deep_research_plan_fill(research_request, preview)
    worker_topics = list(fill.get("deep_research", {}).get("worker_topics") or [])
    signature = deep_research_plan_signature(research_request, worker_topics=worker_topics)
    plan = await get_plan_mode_service().ensure_awaiting_plan_from_fill(
        agent_id=tool_request.context.agent_id,
        intent_type="long_task",
        signature=signature,
        fill=fill,
        original_request=research_request.question,
        source="tool_runtime",
        tenant_id=_parse_uuid_or_none(tool_request.context.tenant_id),
        requested_by_user_id=None,
        metadata_json={
            "deep_research_plan": True,
            "deep_research_tool": tool_request.tool_name,
            "requested_output_format": fill.get("deep_research", {}).get("output_format"),
        },
    )
    return _needs_plan_payload(preview, plan=plan)


def _needs_plan_payload(preview: dict[str, Any], *, plan: Any | None = None) -> dict[str, Any]:
    payload = {
        "ok": False,
        "status": "needs_plan",
        "summary": (
            "Confirm the research plan before running. Relay the clarifying_questions and the proposed plan / "
            "worker_topics to the user; collect their answers and explicit approval first. Preference memory may "
            "prefill parameters but is NOT approval; never self-confirm on the user's behalf."
        ),
        **preview,
        "next_action": (
            "STOP and present the clarifying_questions + proposed worker_topics to the user, then WAIT for their "
            "reply. Only call this tool again with plan_confirmed=true AFTER the user explicitly approves THIS plan "
            "in a new message. Preference memory may prefill depth/scope/language, but it is NOT approval; never "
            "self-confirm or set plan_confirmed=true in the same turn the plan was returned."
        ),
    }
    if plan is not None:
        plan_json = plan.plan_json or {}
        deep_research_plan = plan_json.get("deep_research") if isinstance(plan_json, dict) else {}
        if isinstance(deep_research_plan, dict):
            payload["worker_topics"] = deep_research_plan.get("worker_topics") or payload.get("worker_topics") or []
            payload["clarifying_questions"] = (
                deep_research_plan.get("clarifying_questions") or payload.get("clarifying_questions") or []
            )
        payload.pop("plan", None)
        plan_markdown = _plan_markdown_for_payload(plan)
        payload.update(
            {
                "plan_id": str(plan.id),
                "plan_version": plan.plan_version,
                "plan_hash": plan.plan_hash,
                "plan_markdown_path": plan.plan_markdown_path,
                "plan_markdown": plan_markdown,
                "summary": (
                    "A Deep Research Markdown plan has been created. Show the plan_markdown / plan_markdown_path "
                    "to the user and wait for explicit confirmation before starting the research run."
                ),
            }
        )
    return payload


def _plan_markdown_for_payload(plan: Any) -> str:
    path_value = str(getattr(plan, "plan_markdown_path", "") or "").strip()
    if path_value:
        path = Path(path_value)
        try:
            if path.exists() and path.is_file():
                text = path.read_text(encoding="utf-8").strip()
                if text:
                    return text
        except OSError:
            logger.warning("[DeepResearch] failed to read plan markdown at %s", path_value)

    plan_json = getattr(plan, "plan_json", None) or {}
    if not isinstance(plan_json, dict):
        return ""
    try:
        from app.services import plan_mode_core

        created_at = getattr(plan, "created_at", None)
        return plan_mode_core.render_plan_markdown(
            plan_id=getattr(plan, "id", ""),
            agent_id=getattr(plan, "agent_id", ""),
            tenant_id=getattr(plan, "tenant_id", None),
            status=str(getattr(plan, "status", "awaiting_confirmation") or "awaiting_confirmation"),
            plan_version=int(getattr(plan, "plan_version", 1) or 1),
            plan_hash=str(getattr(plan, "plan_hash", "") or ""),
            intent_type=str(getattr(plan, "intent_type", plan_json.get("intent_type") or "long_task") or "long_task"),
            created_at=created_at.isoformat() if created_at else "",
            plan_json=plan_json,
            confirmed_by=getattr(plan, "confirmed_by_user_id", None),
            confirmed_at=None,
        ).strip()
    except Exception:
        logger.warning("[DeepResearch] failed to render fallback plan markdown for plan=%s", getattr(plan, "id", None))
        return str(plan_json.get("plan_markdown") or plan_json.get("objective") or plan_json.get("title") or "").strip()


@tool(
    ToolMeta(
        name="deep_research_run",
        description=(
            "Run a source-ledger-backed deep research workflow synchronously for quick or standard scopes. "
            "Produces report.md, sources.jsonl, claims.jsonl, steps.jsonl, and final.json artifacts."
        ),
        parameters=_schema(_REQUEST_PROPERTIES, ["question"]),
        category="deep_research_pack",
        display_name="Deep Research Run",
        icon="🔎",
        is_default=False,
        governance="sensitive",
        pack="deep_research_pack",
        adapter="request",
    )
)
async def deep_research_run(request: ToolExecutionRequest) -> str:
    try:
        research_request = ResearchRequest.from_arguments(request.arguments)
    except ValueError as exc:
        return _json({"ok": False, "error": str(exc)})

    if research_request.depth in {"full", "flagship"}:
        return _json(
            {
                "ok": False,
                "status": "route_to_async",
                "error": "async_required",
                "summary": (
                    "Full or flagship Deep Research is a long-running workflow. "
                    "Use deep_research_start instead of deep_research_run."
                ),
                "recommended_tool": "deep_research_start",
                "next_action": (
                    "Call deep_research_start with the same arguments. "
                    "Do not create triggers to poll this task; use deep_research_check "
                    "or the RuntimeTask/artifact UI."
                ),
            }
        )

    if not research_request.plan_confirmed:
        preview = await _plan_preview(research_request)
        return _json(await _materialize_plan_card_payload(request, research_request, preview))

    # DR-6b: the workflow runtime is the ONE Deep Research path. quick/standard
    # depth runs synchronously to completion (export + workspace packet happen
    # inside the launch on completion).
    payload = await start_deep_research_workflow_run(
        request=research_request,
        agent_id=request.context.agent_id,
        user_id=request.context.user_id,
        workspace=request.context.workspace,
    )
    return _json({"ok": payload.get("status") == "completed", **payload})


@tool(
    ToolMeta(
        name="deep_research_start",
        description=(
            "Start a long-running source-ledger-backed deep research workflow. "
            "Creates a RuntimeTask and writes resumable artifacts under runtime_artifacts/long_tasks."
        ),
        parameters=_schema(_REQUEST_PROPERTIES, ["question"]),
        category="deep_research_pack",
        display_name="Deep Research Start",
        icon="🔎",
        is_default=False,
        governance="sensitive",
        # Plan Mode §9.2 bridge: register as plan-governed but keep this tool's
        # own plan_confirmed gate (below) — the service gate must NOT double-block
        # it. Value mirrors app.tools.plan_gate_registry.BRIDGE_SELF.
        plan_gate_action_kind="bridge:self",
        pack="deep_research_pack",
        adapter="request",
    )
)
async def deep_research_start(request: ToolExecutionRequest) -> str:
    try:
        research_request = ResearchRequest.from_arguments(request.arguments)
    except ValueError as exc:
        return _json({"ok": False, "error": str(exc)})

    if not research_request.plan_confirmed:
        preview = await _plan_preview(research_request)
        return _json(await _materialize_plan_card_payload(request, research_request, preview))

    dedup_key = (str(request.context.agent_id), research_request.question.casefold())
    existing_task = _INFLIGHT_DEEP_RESEARCH.get(dedup_key)
    if existing_task:
        return _json(
            {
                "ok": True,
                "task_id": existing_task,
                "status": "running",
                "deduped": True,
                "next_action": (
                    f"A deep research run for this question is already in progress (task {existing_task}). "
                    "Use deep_research_check instead of starting another."
                ),
            }
        )

    # DR-6b single path: pre-generate the workflow run id, schedule the run in
    # the background, and hand the caller a checkable task id immediately —
    # the same async contract the retired RuntimeTask(deep_research) path had.
    run_id = uuid.uuid4()
    _INFLIGHT_DEEP_RESEARCH[dedup_key] = str(run_id)
    _schedule_deep_research_workflow_background(
        research_request=research_request,
        run_id=run_id,
        agent_id=request.context.agent_id,
        user_id=request.context.user_id,
        workspace=request.context.workspace,
    )
    from app.services.deep_research.leaf_presets import run_artifact_dir

    artifact_root = run_artifact_dir(request.context.agent_id, run_id)
    workspace_artifact_dir = _workspace_export_dir(request.context.workspace, run_id.hex)
    return _json(
        {
            "ok": True,
            "task_id": str(run_id),
            "workflow_run_id": str(run_id),
            "status": "running",
            "artifact_dir": str(artifact_root),
            "workspace_artifact_dir": _relative(request.context.workspace, workspace_artifact_dir),
            "next_action": (
                "This research runs asynchronously in the background (usually several minutes). "
                f"Call deep_research_check with task_id {run_id} at most once to confirm it "
                "started; if it is still running, stop here and tell the user the research is "
                "running in the background and they can check back shortly. Do not repeatedly poll "
                "the same running task in one turn: it will not speed it up and will trip the loop "
                "guard. Do not create triggers to poll this task; the RuntimeTask/artifact UI tracks "
                "progress."
            ),
        }
    )


@tool(
    ToolMeta(
        name="deep_research_check",
        description="Check a deep research RuntimeTask, including progress, partial report, source count, claim count, and gaps.",
        parameters=_schema({"task_id": {"type": "string"}}, ["task_id"]),
        category="deep_research_pack",
        display_name="Deep Research Check",
        icon="🔎",
        is_default=False,
        read_only=True,
        parallel_safe=True,
        governance="safe",
        pack="deep_research_pack",
        adapter="request",
    )
)
async def deep_research_check(request: ToolExecutionRequest) -> str:
    task_id = str(request.arguments.get("task_id") or "").strip()
    if not task_id:
        return _json({"ok": False, "error": "task_id_required"})
    record = await get_runtime_task_record(task_id)
    if record and record.get("parent_agent_id") not in {None, str(request.context.agent_id)}:
        return _json({"ok": False, "error": "forbidden"})
    if record and record.get("task_type") == "workflow":
        # DR-4: workflow-shaped Deep Research run — artifacts live under the
        # run-scoped workflow root, progress mirrors the run status.
        artifact_dir = _resolve_deep_research_artifact_dir(
            request.context.workspace,
            task_id,
            agent_id=request.context.agent_id,
            prefer_workflow=True,
        )
        payload = _read_artifact_dir_payload(request.context.workspace, task_id, artifact_dir)
        if not payload.get("status"):
            payload["status"] = record.get("status")
        return _json({"ok": True, "task": record, **payload})
    payload = _read_deep_research_artifact(request.context.workspace, task_id)
    return _json({"ok": True, "task": record, **payload})


@tool(
    ToolMeta(
        name="deep_research_cancel",
        description="Cancel a running deep research RuntimeTask that belongs to the current agent.",
        parameters=_schema({"task_id": {"type": "string"}, "reason": {"type": "string"}}, ["task_id"]),
        category="deep_research_pack",
        display_name="Deep Research Cancel",
        icon="🔎",
        is_default=False,
        governance="sensitive",
        pack="deep_research_pack",
        adapter="request",
    )
)
async def deep_research_cancel(request: ToolExecutionRequest) -> str:
    task_id = str(request.arguments.get("task_id") or "").strip()
    if not task_id:
        return _json({"ok": False, "error": "task_id_required"})
    record = await get_runtime_task_record(task_id)
    if record and record.get("parent_agent_id") not in {None, str(request.context.agent_id)}:
        return _json({"ok": False, "error": "forbidden"})
    await update_runtime_task_record(
        task_id,
        status="killed",
        result_summary=str(request.arguments.get("reason") or "Deep research cancelled."),
        metadata_json={"cancel_reason": str(request.arguments.get("reason") or "cancelled")},
    )
    parsed = _parse_uuid(task_id)
    if parsed:
        await record_long_task_progress(
            agent_id=request.context.agent_id,
            runtime_task_id=parsed,
            status="killed",
            delta="Deep research task was cancelled.",
            blocked_reason=str(request.arguments.get("reason") or "cancelled"),
        )
    return _json({"ok": True, "task_id": task_id, "status": "killed"})


@tool(
    ToolMeta(
        name="deep_research_export",
        description="Export a completed or partial deep research artifact as markdown, json, html, docx, xlsx, or pptx.",
        parameters=_schema(
            {
                "task_id": {"type": "string"},
                "format": {"type": "string", "enum": ["markdown", "json", "html", "docx", "xlsx", "pptx"]},
            },
            ["task_id"],
        ),
        category="deep_research_pack",
        display_name="Deep Research Export",
        icon="🔎",
        is_default=False,
        governance="sensitive",
        pack="deep_research_pack",
        adapter="request",
    )
)
async def deep_research_export(request: ToolExecutionRequest) -> str:
    task_id = str(request.arguments.get("task_id") or "").strip()
    export_format = normalize_deep_research_output_format(str(request.arguments.get("format") or "markdown").strip())
    record = await get_runtime_task_record(task_id)
    if record and record.get("parent_agent_id") not in {None, str(request.context.agent_id)}:
        return _json({"ok": False, "error": "forbidden"})
    artifact_dir = _resolve_deep_research_artifact_dir(
        request.context.workspace,
        task_id,
        agent_id=request.context.agent_id,
        prefer_workflow=bool(record and record.get("task_type") == "workflow"),
    )
    final_path = artifact_dir / "final.json"
    report_path = artifact_dir / "report.md"
    artifact_target = _materialize_requested_output_format(request.context.workspace, artifact_dir, export_format)
    if artifact_target is None:
        artifact_target = final_path if export_format == "json" else report_path
    file_name = artifact_target.name
    if not artifact_target.exists():
        return _json({"ok": False, "error": "artifact_not_found", "task_id": task_id})

    workspace_dir = _publish_workspace_packet(request.context.workspace, task_id, artifact_dir)
    target = workspace_dir / file_name
    return _json(
        {
            "ok": True,
            "task_id": task_id,
            "format": export_format,
            "path": _relative(request.context.workspace, target),
            "workspace_artifact_dir": _relative(request.context.workspace, workspace_dir),
            "artifact_path": _relative(request.context.workspace, artifact_target),
        }
    )


async def start_deep_research_background_run(
    *,
    request: ResearchRequest,
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    workspace: Path,
    plan_id: uuid.UUID | str | None = None,
) -> dict[str, Any]:
    """Plan-handoff entry (confirmed plan → background workflow run).

    DR-6b single path: pre-generate the run id, schedule the workflow in the
    background, return a checkable id immediately — confirmation must never
    block on a multi-minute research run.
    """
    run_id = uuid.uuid4()
    _INFLIGHT_DEEP_RESEARCH[(str(agent_id), request.question.casefold())] = str(run_id)
    _schedule_deep_research_workflow_background(
        research_request=request,
        run_id=run_id,
        agent_id=agent_id,
        user_id=user_id,
        workspace=workspace,
        plan_id=plan_id,
    )
    from app.services.deep_research.leaf_presets import run_artifact_dir

    return {
        "created_workflow_run_id": str(run_id),
        "task_id": str(run_id),
        "status": "running",
        "artifact_dir": str(run_artifact_dir(agent_id, run_id)),
        "workspace_artifact_dir": _relative(workspace, _workspace_export_dir(workspace, run_id.hex)),
        "output_format": normalize_deep_research_output_format(request.output_format),
    }


def _schedule_deep_research_workflow_background(
    *,
    research_request: ResearchRequest,
    run_id: uuid.UUID,
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    workspace: Path,
    plan_id: uuid.UUID | str | None = None,
) -> None:
    """Run the deep_research.v1 workflow in the background.

    The workflow runtime owns the RuntimeTask row, journal, status and
    completion-edge artifact landing; this runner only annotates the result
    payload onto the task metadata and clears the in-flight dedup slot."""

    async def runner() -> None:
        try:
            payload = await start_deep_research_workflow_run(
                request=research_request,
                agent_id=agent_id,
                user_id=user_id,
                workspace=workspace,
                plan_id=plan_id,
                run_id=run_id,
            )
            await update_runtime_task_record(
                run_id.hex,
                result_summary=f"Deep research workflow {payload.get('status')}",
                metadata_json={"deep_research_result": payload},
            )
        except Exception as exc:
            logger.warning("[DeepResearch] background workflow run %s failed: %s", run_id, exc)
            try:
                await update_runtime_task_record(
                    run_id.hex, result_summary=f"Deep research failed: {type(exc).__name__}: {exc}"
                )
            except Exception:
                logger.warning("[DeepResearch] failed to annotate failure on run %s", run_id, exc_info=True)
        finally:
            _INFLIGHT_DEEP_RESEARCH.pop((str(agent_id), research_request.question.casefold()), None)

    asyncio.create_task(runner())


def _read_deep_research_artifact(workspace: Path, task_id: str) -> dict[str, Any]:
    return _read_artifact_dir_payload(workspace, task_id, _resolve_deep_research_artifact_dir(workspace, task_id))


def _read_artifact_dir_payload(workspace: Path, task_id: str, artifact_dir: Path) -> dict[str, Any]:
    """Artifact-dir-injected progress payload — shared by the legacy
    long_tasks layout and the workflow run-scoped layout (DR-4)."""
    if (artifact_dir / "report.md").exists() or (artifact_dir / "final.json").exists():
        _publish_workspace_packet(workspace, task_id, artifact_dir)
    workspace_dir = _workspace_export_dir(workspace, task_id)
    workspace_report_path = workspace_dir / "report.md"
    workspace_sources_path = workspace_dir / "sources.jsonl"
    workspace_claims_path = workspace_dir / "claims.jsonl"
    workspace_final_path = workspace_dir / "final.json"
    workspace_source_notes_path = workspace_dir / "source_notes.jsonl"
    workspace_lane_summaries_path = workspace_dir / "lane_summaries.jsonl"
    final = _load_json(artifact_dir / "final.json") or {}
    return {
        "artifact_dir": _relative(workspace, artifact_dir),
        "workspace_artifact_dir": _relative(workspace, workspace_dir) if workspace_dir.exists() else None,
        "status": final.get("status"),
        "summary": final.get("summary"),
        "report_path": _relative(workspace, workspace_report_path)
        if workspace_report_path.exists()
        else (_relative(workspace, artifact_dir / "report.md") if (artifact_dir / "report.md").exists() else None),
        "sources_path": _relative(workspace, workspace_sources_path)
        if workspace_sources_path.exists()
        else (
            _relative(workspace, artifact_dir / "sources.jsonl") if (artifact_dir / "sources.jsonl").exists() else None
        ),
        "claims_path": _relative(workspace, workspace_claims_path)
        if workspace_claims_path.exists()
        else (
            _relative(workspace, artifact_dir / "claims.jsonl") if (artifact_dir / "claims.jsonl").exists() else None
        ),
        "source_notes_path": _relative(workspace, workspace_source_notes_path)
        if workspace_source_notes_path.exists()
        else (
            _relative(workspace, artifact_dir / "source_notes.jsonl")
            if (artifact_dir / "source_notes.jsonl").exists()
            else None
        ),
        "lane_summaries_path": _relative(workspace, workspace_lane_summaries_path)
        if workspace_lane_summaries_path.exists()
        else (
            _relative(workspace, artifact_dir / "lane_summaries.jsonl")
            if (artifact_dir / "lane_summaries.jsonl").exists()
            else None
        ),
        "final_path": _relative(workspace, workspace_final_path)
        if workspace_final_path.exists()
        else (_relative(workspace, artifact_dir / "final.json") if (artifact_dir / "final.json").exists() else None),
        "artifact_report_path": _relative(workspace, artifact_dir / "report.md")
        if (artifact_dir / "report.md").exists()
        else None,
        "source_count": final.get("source_count") or _jsonl_count(artifact_dir / "sources.jsonl"),
        "claim_count": final.get("claim_count") or _jsonl_count(artifact_dir / "claims.jsonl"),
        "quality_gates": final.get("quality_gates") or {},
        "gaps": final.get("gaps") or [],
        "partial_report": (artifact_dir / "report.md").read_text(encoding="utf-8")[:4000]
        if (artifact_dir / "report.md").exists()
        else "",
    }


def _deep_research_dir(workspace: Path, task_id: uuid.UUID | str) -> Path:
    return workspace / "runtime_artifacts" / "long_tasks" / str(task_id).replace("-", "") / "deep_research"


def _workflow_deep_research_dir(workspace: Path, task_id: uuid.UUID | str) -> Path:
    return workspace / "runtime_artifacts" / "workflow_runs" / str(task_id) / "deep_research"


_ARTIFACT_SIGNAL_FILES = (
    "final.json",
    "report.md",
    "sources.jsonl",
    "claims.jsonl",
    "source_notes.jsonl",
    "lane_summaries.jsonl",
)


def _artifact_dir_has_material(path: Path) -> bool:
    return any((path / file_name).exists() for file_name in _ARTIFACT_SIGNAL_FILES)


def _run_artifact_dir_for_agent(agent_id: uuid.UUID | str | None, task_id: uuid.UUID | str) -> Path | None:
    if not agent_id:
        return None
    try:
        from app.services.deep_research.leaf_presets import run_artifact_dir

        return run_artifact_dir(agent_id, task_id)
    except Exception:
        logger.warning("[DeepResearch] failed to resolve workflow artifact dir for agent=%s run=%s", agent_id, task_id)
        return None


def _resolve_deep_research_artifact_dir(
    workspace: Path,
    task_id: uuid.UUID | str,
    *,
    agent_id: uuid.UUID | str | None = None,
    prefer_workflow: bool = False,
) -> Path:
    """Resolve Deep Research artifacts across the workflow and legacy layouts.

    Workflow runs write under ``runtime_artifacts/workflow_runs/{run}/deep_research``.
    Legacy synchronous runs wrote under ``runtime_artifacts/long_tasks/{run}/deep_research``.
    """
    candidates: list[Path | None] = []
    if prefer_workflow:
        candidates.append(_run_artifact_dir_for_agent(agent_id, task_id))
        candidates.append(_workflow_deep_research_dir(workspace, task_id))
        candidates.append(_deep_research_dir(workspace, task_id))
    else:
        candidates.append(_workflow_deep_research_dir(workspace, task_id))
        candidates.append(_deep_research_dir(workspace, task_id))
        candidates.append(_run_artifact_dir_for_agent(agent_id, task_id))

    seen: set[str] = set()
    resolved_candidates: list[Path] = []
    for candidate in candidates:
        if candidate is None:
            continue
        key = str(candidate)
        if key in seen:
            continue
        seen.add(key)
        resolved_candidates.append(candidate)
        if _artifact_dir_has_material(candidate):
            return candidate
    return resolved_candidates[0] if resolved_candidates else _deep_research_dir(workspace, task_id)


def _workspace_export_dir(workspace: Path, run_id: uuid.UUID | str) -> Path:
    safe_id = str(run_id).strip() or "latest"
    return workspace / "workspace" / "deep_research_reports" / safe_id


def _publish_workspace_packet(workspace: Path, run_id: uuid.UUID | str, artifact_dir: Path) -> Path:
    workspace_dir = _workspace_export_dir(workspace, run_id)
    workspace_dir.mkdir(parents=True, exist_ok=True)
    for file_name in (
        "request.json",
        "plan.json",
        "steps.jsonl",
        "sources.jsonl",
        "claims.jsonl",
        "evaluation.jsonl",
        "source_notes.jsonl",
        "lane_summaries.jsonl",
        "devils_advocate.jsonl",
        "report.md",
        "report.html",
        "report.docx",
        "report.xlsx",
        "report.pptx",
        "final.json",
    ):
        source = artifact_dir / file_name
        if source.exists() and source.is_file():
            shutil.copyfile(source, workspace_dir / file_name)
    return workspace_dir


def _materialize_requested_output_format(workspace: Path, artifact_dir: Path, output_format: str) -> Path | None:
    return compose_deep_research_artifact(workspace, artifact_dir, output_format)


def _save_office_bytes(workspace: Path, target: Path, content: bytes, *, reason: str) -> None:
    rel_path = _relative(workspace, target)
    OfficeDocumentService(workspace).atomic_save_bytes(
        rel_path,
        content,
        reason=reason,
        require_no_active_editor=True,
    )


def _markdown_to_html_document(markdown: str) -> str:
    body = "\n".join(_markdown_line_to_html(line) for line in markdown.splitlines())
    return (
        "<!doctype html>\n"
        '<html><head><meta charset="utf-8"><title>Deep Research Report</title>'
        "<style>body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;line-height:1.55;"
        "max-width:920px;margin:40px auto;padding:0 24px;color:#111827}"
        "code,pre{font-family:ui-monospace,SFMono-Regular,Menlo,monospace}"
        "blockquote{border-left:3px solid #d1d5db;margin-left:0;padding-left:14px;color:#4b5563}"
        "</style></head><body>\n"
        f"{body}\n"
        "</body></html>\n"
    )


def _markdown_line_to_html(line: str) -> str:
    stripped = line.strip()
    if not stripped:
        return ""
    if stripped.startswith("### "):
        return f"<h3>{html.escape(stripped[4:])}</h3>"
    if stripped.startswith("## "):
        return f"<h2>{html.escape(stripped[3:])}</h2>"
    if stripped.startswith("# "):
        return f"<h1>{html.escape(stripped[2:])}</h1>"
    if stripped.startswith("- "):
        return f"<ul><li>{html.escape(stripped[2:])}</li></ul>"
    if stripped.startswith("> "):
        return f"<blockquote>{html.escape(stripped[2:])}</blockquote>"
    return f"<p>{html.escape(stripped)}</p>"


def _markdown_to_docx_bytes(markdown: str) -> bytes:
    from docx import Document

    document = Document()
    for line in markdown.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            document.add_heading(stripped[2:].strip(), level=1)
        elif stripped.startswith("## "):
            document.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("### "):
            document.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith("- "):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
        else:
            document.add_paragraph(stripped)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _markdown_to_pptx_bytes(markdown: str) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    title = "Deep Research Report"
    sections: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_body: list[str] = []
    for line in markdown.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            if current_title:
                sections.append((current_title, current_body))
            current_title = stripped.lstrip("#").strip()
            current_body = []
            if title == "Deep Research Report":
                title = current_title or title
        elif current_title:
            current_body.append(stripped[2:].strip() if stripped.startswith("- ") else stripped)
    if current_title:
        sections.append((current_title, current_body))
    if not sections:
        sections = [(title, [line.strip() for line in markdown.splitlines() if line.strip()][:8])]

    for index, (section_title, body_lines) in enumerate(sections[:12]):
        layout = presentation.slide_layouts[0] if index == 0 else presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = section_title[:120]
        placeholders = [shape for shape in slide.placeholders if shape.placeholder_format.idx != 0]
        if placeholders:
            placeholders[0].text = "\n".join(body_lines[:8])
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _output_file_name(output_format: str) -> str:
    output_format = normalize_deep_research_output_format(output_format)
    return {
        "markdown": "report.md",
        "json": "final.json",
        "html": "report.html",
        "docx": "report.docx",
        "xlsx": "report.xlsx",
        "pptx": "report.pptx",
    }[output_format]


def _relative(workspace: Path, path: Path) -> str:
    try:
        return path.resolve().relative_to(workspace.resolve()).as_posix()
    except ValueError:
        return path.as_posix()


def _load_json(path: Path) -> dict[str, Any] | None:
    if not path.exists():
        return None
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    return payload if isinstance(payload, dict) else None


def _jsonl_count(path: Path) -> int:
    if not path.exists():
        return 0
    return len([line for line in path.read_text(encoding="utf-8").splitlines() if line.strip()])


def _parse_uuid(value: str) -> uuid.UUID | None:
    try:
        return uuid.UUID(str(value))
    except (TypeError, ValueError):
        return None


def _parse_uuid_or_none(value: str | uuid.UUID | None) -> uuid.UUID | None:
    if isinstance(value, uuid.UUID):
        return value
    if not value:
        return None
    return _parse_uuid(str(value))


def _escape_html(value: str) -> str:
    return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _json(payload: dict[str, Any]) -> str:
    return json.dumps(to_jsonable(payload), ensure_ascii=False, default=str)


# ──────────────────────────────────────────────────────────────────────────
# Tier 3-4: streaming generator for SSE-style consumption.
#
# Yields incremental events as new lines appear in steps.jsonl /
# claims.jsonl / source_notes.jsonl / lane_summaries.jsonl / reflection.jsonl
# and finally a `report` event with the partial markdown plus a `final`
# event when final.json appears. Idempotent across reconnects: callers can
# pass `after_step_index` / `after_claim_index` cursors to resume.
#
# Designed to be wrapped by an SSE FastAPI route (deferred to API layer);
# pure async-generator contract here so the kernel and tests can drive it.
# ──────────────────────────────────────────────────────────────────────────

_STREAM_FILES: tuple[tuple[str, str], ...] = (
    ("step", "steps.jsonl"),
    ("claim", "claims.jsonl"),
    ("source_note", "source_notes.jsonl"),
    ("lane_summary", "lane_summaries.jsonl"),
    ("reflection", "reflection.jsonl"),
    ("devils_advocate", "devils_advocate.jsonl"),
    ("controller_trace", "controller_trace.jsonl"),
)


async def stream_deep_research_artifacts(
    workspace: Path,
    task_id: str,
    *,
    poll_interval_seconds: float = 0.5,
    cursors: dict[str, int] | None = None,
    deadline_seconds: float | None = None,
    _now: Any = None,
) -> AsyncIterator[dict[str, Any]]:
    """Async-generator that streams Deep Research artifact deltas as discrete events.

    Each yielded event is a JSON-serialisable dict with at minimum:
      - `event`: one of step/claim/source_note/lane_summary/reflection/
        controller_trace/report/final/heartbeat
      - `task_id`
      - `timestamp` (ISO8601)
      - `payload`: the decoded record (for line-oriented files) or text body

    The generator terminates when:
      - `final.json` has been emitted, OR
      - `deadline_seconds` is exceeded since the first call, OR
      - the caller closes / cancels the iterator.

    `cursors` lets a reconnecting caller resume past previously-seen records.
    """
    import time as _time_mod
    from datetime import datetime, timezone

    artifact_dir = _resolve_deep_research_artifact_dir(workspace, task_id)
    cursors = dict(cursors or {})
    started_at = (_now or _time_mod.monotonic)()
    final_emitted = False
    last_report_text = ""

    def _ts() -> str:
        return datetime.now(timezone.utc).isoformat()

    while True:
        emitted_any = False

        for event_name, filename in _STREAM_FILES:
            path = artifact_dir / filename
            if not path.exists():
                continue
            try:
                lines = [line for line in path.read_text("utf-8").splitlines() if line.strip()]
            except OSError:
                continue
            cursor = cursors.get(event_name, 0)
            if cursor >= len(lines):
                continue
            for line in lines[cursor:]:
                try:
                    payload = json.loads(line)
                except json.JSONDecodeError:
                    payload = {"raw": line}
                yield {
                    "event": event_name,
                    "task_id": task_id,
                    "timestamp": _ts(),
                    "payload": payload,
                }
                emitted_any = True
            cursors[event_name] = len(lines)

        report_path = artifact_dir / "report.md"
        if report_path.exists():
            try:
                text = report_path.read_text("utf-8")
            except OSError:
                text = ""
            if text != last_report_text:
                yield {
                    "event": "report",
                    "task_id": task_id,
                    "timestamp": _ts(),
                    "payload": {"markdown": text, "chars": len(text)},
                }
                last_report_text = text
                emitted_any = True

        final_path = artifact_dir / "final.json"
        if final_path.exists() and not final_emitted:
            try:
                final_payload = json.loads(final_path.read_text("utf-8"))
            except (OSError, json.JSONDecodeError):
                final_payload = {}
            yield {
                "event": "final",
                "task_id": task_id,
                "timestamp": _ts(),
                "payload": final_payload,
            }
            final_emitted = True
            return

        if not emitted_any:
            yield {
                "event": "heartbeat",
                "task_id": task_id,
                "timestamp": _ts(),
                "payload": {"cursors": dict(cursors)},
            }

        if deadline_seconds is not None:
            elapsed = (_now or _time_mod.monotonic)() - started_at
            if elapsed >= deadline_seconds:
                return

        await asyncio.sleep(poll_interval_seconds)
