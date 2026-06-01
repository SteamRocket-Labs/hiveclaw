from __future__ import annotations

import json
from pathlib import Path


def _write_sample_dossier(artifact_dir: Path) -> None:
    artifact_dir.mkdir(parents=True, exist_ok=True)
    (artifact_dir / "report.md").write_text(
        """# RWA Custody Research

## Executive Thesis

The evidence supports a cautious launch: adoption is rising, but custody and reporting controls
remain the gating factor.

## Key Findings

- Primary disclosures show 35% growth across 12 jurisdictions.
- The custody workflow has 18 controls and 7 reporting checkpoints.

## Contradictions And Gaps

- Independent audit evidence is still thin.

## Strategic Implications

Prioritize regulated custody integrations before broad distribution.

## Source Ledger

- `src_a` Official disclosure
- `src_b` Regulator guidance
""",
        encoding="utf-8",
    )
    (artifact_dir / "final.json").write_text(
        json.dumps(
            {
                "status": "completed",
                "summary": "Completed with 2 sources and 3 claims.",
                "quality_gates": {"synthesis": "passed"},
                "gaps": ["Independent audit evidence is still thin."],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    (artifact_dir / "sources.jsonl").write_text(
        "\n".join(
            [
                json.dumps(
                    {
                        "source_id": "src_a",
                        "title": "Official disclosure",
                        "publisher": "Issuer",
                        "url": "https://example.com/a",
                        "source_type": "primary",
                        "evidence_tier": "tier1",
                        "evidence_grade": "A",
                    }
                ),
                json.dumps(
                    {
                        "source_id": "src_b",
                        "title": "Regulator guidance",
                        "publisher": "Regulator",
                        "url": "https://example.com/b",
                        "source_type": "regulatory",
                        "evidence_tier": "tier1",
                        "evidence_grade": "A",
                    }
                ),
            ]
        ),
        encoding="utf-8",
    )
    (artifact_dir / "claims.jsonl").write_text(
        json.dumps(
            {
                "claim_id": "claim_a",
                "text": "Adoption grew 35% across 12 jurisdictions.",
                "status": "verified",
                "source_ids": ["src_a"],
                "evidence": "Official disclosure",
            }
        ),
        encoding="utf-8",
    )
    (artifact_dir / "lane_summaries.jsonl").write_text(
        json.dumps(
            {
                "lane_id": "official",
                "label": "Official evidence",
                "source_count": 2,
                "evidence_strength": "strong",
                "key_findings": ["35% growth", "18 controls"],
            }
        ),
        encoding="utf-8",
    )


def test_artifact_composer_supports_all_deep_research_formats(tmp_path):
    from app.services.deep_research.artifact_composer import SUPPORTED_DEEP_RESEARCH_OUTPUT_FORMATS

    assert SUPPORTED_DEEP_RESEARCH_OUTPUT_FORMATS == {"markdown", "json", "html", "docx", "xlsx", "pptx"}


def test_artifact_composer_builds_xlsx_evidence_workbook(tmp_path):
    from openpyxl import load_workbook

    from app.services.deep_research.artifact_composer import compose_deep_research_artifact

    artifact_dir = tmp_path / "runtime_artifacts" / "long_tasks" / "task" / "deep_research"
    _write_sample_dossier(artifact_dir)

    target = compose_deep_research_artifact(tmp_path, artifact_dir, "xlsx")

    assert target.name == "report.xlsx"
    workbook = load_workbook(target)
    assert {"Summary", "Sources", "Claims", "Lane Coverage", "Quality Gates"}.issubset(workbook.sheetnames)
    assert workbook["Claims"]["B2"].value == "Adoption grew 35% across 12 jurisdictions."


def test_artifact_composer_builds_docx_memo_not_plain_markdown_dump(tmp_path):
    from docx import Document

    from app.services.deep_research.artifact_composer import compose_deep_research_artifact

    artifact_dir = tmp_path / "runtime_artifacts" / "long_tasks" / "task" / "deep_research"
    _write_sample_dossier(artifact_dir)

    target = compose_deep_research_artifact(tmp_path, artifact_dir, "docx")

    document = Document(target)
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Executive Summary" in text
    assert "Evidence Appendix" in text
    assert "# RWA Custody Research" not in text


def test_artifact_composer_builds_pptx_decision_deck(tmp_path):
    from pptx import Presentation

    from app.services.deep_research.artifact_composer import compose_deep_research_artifact

    artifact_dir = tmp_path / "runtime_artifacts" / "long_tasks" / "task" / "deep_research"
    _write_sample_dossier(artifact_dir)

    target = compose_deep_research_artifact(tmp_path, artifact_dir, "pptx")

    presentation = Presentation(target)
    titles = [slide.shapes.title.text for slide in presentation.slides if slide.shapes.title]
    assert titles[:4] == [
        "Decision Question",
        "Executive Thesis",
        "Evidence Matrix",
        "Risks And Gaps",
    ]
    assert len(titles) <= 8
