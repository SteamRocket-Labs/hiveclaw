from __future__ import annotations

import html as html_lib
import json
import logging
import os
import re
import shutil
import tempfile
from dataclasses import dataclass, replace
from datetime import UTC, datetime
from hashlib import sha256
from pathlib import Path
from typing import Any, Literal

from app.config import get_settings
from app.services.officecli_adapter import OfficeCLIAdapter
from app.services.officecli_adapter import OfficeCLIError


logger = logging.getLogger(__name__)


SUPPORTED_OFFICE_KINDS: frozenset[str] = frozenset({"docx", "xlsx", "pptx"})
SUPPORTED_OFFICE_SUFFIXES: dict[str, str] = {
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
}

OFFICE_PREVIEW_CSP = (
    "sandbox allow-scripts; default-src 'none'; img-src data: blob:; "
    "style-src 'unsafe-inline'; script-src 'unsafe-inline'; font-src data:; "
    "connect-src 'none'; form-action 'none'; base-uri 'none'; object-src 'none'"
)


@dataclass(frozen=True)
class OfficePreviewResult:
    html: str
    preview_mode: Literal["html", "text_fallback"]
    source_sha256: str
    renderer_version: str
    cache_hit: bool
    output_bytes: int
    fallback_reason: str | None = None


class OfficeDocumentError(RuntimeError):
    error_code = "office_document_error"


class OfficeDocumentPathError(OfficeDocumentError):
    error_code = "auth_or_permission"


class OfficeDocumentExistsError(OfficeDocumentError):
    error_code = "already_exists"


class OfficeDocumentNotFoundError(OfficeDocumentError):
    error_code = "not_found"


class OfficePreviewError(OfficeDocumentError):
    error_code = "office_preview_error"


class OfficePreviewUnsupportedTypeError(OfficePreviewError):
    error_code = "office_preview_unsupported_type"


class OfficePreviewTooLargeError(OfficePreviewError):
    error_code = "office_preview_too_large"


class OfficePreviewSourceChangedError(OfficePreviewError):
    error_code = "office_preview_source_changed"


class OfficePreviewUnavailableError(OfficePreviewError):
    error_code = "office_preview_unavailable"


class OfficePreviewMalformedError(OfficePreviewError):
    error_code = "office_preview_malformed"


class OfficeDocumentService:
    """Workspace-scoped document manager for Office artifacts and sidecars."""

    def __init__(
        self,
        workspace: Path,
        *,
        adapter: OfficeCLIAdapter | None = None,
        preview_max_bytes: int | None = None,
    ) -> None:
        settings = get_settings()
        self.workspace = workspace.resolve()
        self.adapter = adapter or OfficeCLIAdapter()
        self.preview_max_bytes = int(
            preview_max_bytes
            if preview_max_bytes is not None
            else getattr(settings, "OFFICECLI_PREVIEW_MAX_BYTES", 25 * 1024 * 1024)
        )
        if self.preview_max_bytes <= 0:
            raise ValueError("Office preview byte limit must be positive")

    def resolve_document_path(self, rel_path: str) -> Path:
        normalized = self._normalize_rel_path(rel_path)
        target = (self.workspace / normalized).resolve()
        self._assert_under_workspace(target)
        return target

    def manifest_path(self, rel_path: str) -> Path:
        normalized = self._normalize_rel_path(rel_path)
        digest = sha256(normalized.encode("utf-8")).hexdigest()
        return self.workspace / ".office_meta" / digest / "manifest.json"

    def create_document(
        self,
        rel_path: str,
        *,
        kind: str,
        template_path: str | None = None,
    ) -> dict[str, Any]:
        target = self.resolve_document_path(rel_path)
        normalized = self._normalize_rel_path(rel_path)
        kind = (kind or "").strip().lower()
        if kind not in SUPPORTED_OFFICE_KINDS:
            raise ValueError(f"Unsupported office document kind: {kind}")
        if target.suffix.lower() != f".{kind}":
            raise ValueError(f"Document path suffix must match kind .{kind}")
        if target.exists():
            raise OfficeDocumentExistsError(f"Document already exists: {normalized}")

        target.parent.mkdir(parents=True, exist_ok=True)
        if template_path:
            template = self.resolve_document_path(template_path)
            if not template.is_file():
                raise OfficeDocumentNotFoundError(f"Template not found: {template_path}")
            self._atomic_copy(template, target)
        else:
            self._create_blank_document(target, kind)

        manifest = self._load_manifest(normalized)
        manifest.update(
            {
                "path": normalized,
                "kind": kind,
                "created_at": self._now(),
                "updated_at": self._now(),
                "current_version": 0,
            }
        )
        self._save_manifest(normalized, manifest)
        return {"path": normalized, "kind": kind, "size": target.stat().st_size}

    def atomic_save_bytes(
        self,
        rel_path: str,
        content: bytes,
        *,
        reason: str,
    ) -> dict[str, Any]:
        normalized = self._normalize_rel_path(rel_path)
        target = self.resolve_document_path(normalized)
        manifest = self._load_manifest(normalized)

        next_version = int(manifest.get("current_version") or 0) + 1
        revision_entry: dict[str, Any] | None = None
        if target.exists():
            revision_entry = self._write_revision(normalized, target, version=next_version, reason=reason)

        target.parent.mkdir(parents=True, exist_ok=True)
        self._atomic_write_bytes(target, content)

        revisions = list(manifest.get("revisions") or [])
        if revision_entry:
            revisions.append(revision_entry)
        manifest.update(
            {
                "path": normalized,
                "kind": self._kind_for_path(target),
                "current_version": next_version,
                "updated_at": self._now(),
                "revisions": revisions,
            }
        )
        self._save_manifest(normalized, manifest)
        return {"path": normalized, "version": next_version, "size": len(content)}

    def run_view(self, rel_path: str, *, mode: str = "outline", page: int | None = None) -> dict[str, Any]:
        target = self._require_existing_file(rel_path)
        return self.adapter.run_view(target, mode=mode, page=page, cwd=self.workspace)

    def run_query(self, rel_path: str, *, selector: str, depth: int | None = None) -> dict[str, Any]:
        target = self._require_existing_file(rel_path)
        options: dict[str, Any] = {"selector": selector}
        if depth is not None:
            options["depth"] = depth
        return self.adapter.run("query", target, options=options, cwd=self.workspace)

    def run_validate(self, rel_path: str) -> dict[str, Any]:
        target = self._require_existing_file(rel_path)
        return self.adapter.run("validate", target, cwd=self.workspace)

    def run_dump(self, rel_path: str) -> dict[str, Any]:
        target = self._require_existing_file(rel_path)
        return self.adapter.run("get", target, cwd=self.workspace)

    def run_apply(
        self,
        rel_path: str,
        *,
        operations: list[dict[str, Any]],
        output_path: str | None = None,
    ) -> dict[str, Any]:
        normalized = self._normalize_rel_path(rel_path)
        target = self._require_existing_file(normalized)

        op_file = self._write_operation(normalized, operations)
        options: dict[str, Any] = {"operations": op_file}
        if output_path:
            options["output"] = self.resolve_document_path(output_path)
        payload = self.adapter.run("batch", target, options=options, cwd=self.workspace)

        manifest = self._load_manifest(normalized)
        manifest["updated_at"] = self._now()
        manifest.setdefault("operations", []).append(
            {
                "file": Path(op_file).name,
                "created_at": self._now(),
                "operation_count": len(operations),
                "output_path": output_path,
            }
        )
        self._save_manifest(normalized, manifest)
        return payload

    def render_preview(self, rel_path: str) -> OfficePreviewResult:
        """Render and cache a workspace Office document preview."""

        normalized = self._normalize_rel_path(rel_path)
        target = self._require_existing_file(normalized)
        cache_dir = self.manifest_path(normalized).parent / "preview"
        return self._render_preview_target(target, cache_dir=cache_dir)

    def render_preview_target(self, target: Path, *, cache_key: str) -> OfficePreviewResult:
        """Render an authorized artifact snapshot already resolved by the API."""

        resolved = target.resolve()
        self._assert_under_workspace(resolved)
        if not resolved.is_file():
            raise OfficeDocumentNotFoundError(f"Office preview source not found: {resolved.name}")
        artifact_key = cache_key.removeprefix("artifact:") if cache_key.startswith("artifact:") else ""
        if artifact_key and re.fullmatch(r"[A-Za-z0-9-]+", artifact_key):
            cache_dir = self.workspace / ".office_meta" / "artifacts" / artifact_key / "preview"
        else:
            digest = sha256(cache_key.encode("utf-8")).hexdigest()
            cache_dir = self.workspace / ".office_meta" / "previews" / digest / "preview"
        return self._render_preview_target(resolved, cache_dir=cache_dir)

    def _render_preview_target(self, target: Path, *, cache_dir: Path) -> OfficePreviewResult:
        if target.suffix.lower() not in SUPPORTED_OFFICE_SUFFIXES:
            raise OfficePreviewUnsupportedTypeError(f"Unsupported Office preview type: {target.suffix or '<none>'}")

        renderer_version = self._renderer_version()
        source_hash = self._sha256_file(target)
        cached = self._load_preview_cache(
            cache_dir,
            source_sha256=source_hash,
            renderer_version=renderer_version,
        )
        if cached is not None and self._sha256_file(target) == source_hash:
            return replace(cached, cache_hit=True)

        for attempt in range(2):
            source_hash = self._sha256_file(target)
            rendered = self._render_preview_once(target, source_hash=source_hash, renderer_version=renderer_version)
            if self._sha256_file(target) == source_hash:
                self._save_preview_cache(cache_dir, rendered)
                return rendered
            logger.warning(
                "Office preview source changed during render: source=%s attempt=%s",
                target.suffix.lower(),
                attempt + 1,
            )
        raise OfficePreviewSourceChangedError("Office preview source changed during both render attempts")

    def _render_preview_once(
        self,
        target: Path,
        *,
        source_hash: str,
        renderer_version: str,
    ) -> OfficePreviewResult:
        try:
            payload = self.adapter.run_view(target, mode="html", cwd=self.workspace)
            raw_html = self._payload_text(payload, mode="html")
            rendered_html = self._harden_html(raw_html)
            output_bytes = self._checked_output_size(rendered_html)
            return OfficePreviewResult(
                html=rendered_html,
                preview_mode="html",
                source_sha256=source_hash,
                renderer_version=renderer_version,
                cache_hit=False,
                output_bytes=output_bytes,
            )
        except OfficePreviewTooLargeError:
            raise
        except (OfficeCLIError, OfficePreviewMalformedError) as html_error:
            fallback_reason = getattr(html_error, "error_code", type(html_error).__name__)

        try:
            payload = self.adapter.run_view(target, mode="text", cwd=self.workspace)
            text = self._payload_text(payload, mode="text")
        except (OfficeCLIError, OfficePreviewMalformedError) as text_error:
            raise OfficePreviewUnavailableError(
                f"OfficeCLI HTML and text preview modes failed ({type(text_error).__name__})"
            ) from text_error

        rendered_html = self._text_fallback_html(text)
        output_bytes = self._checked_output_size(rendered_html)
        return OfficePreviewResult(
            html=rendered_html,
            preview_mode="text_fallback",
            source_sha256=source_hash,
            renderer_version=renderer_version,
            cache_hit=False,
            output_bytes=output_bytes,
            fallback_reason=str(fallback_reason),
        )

    def _renderer_version(self) -> str:
        try:
            return self.adapter.version()
        except OfficeCLIError as exc:
            logger.warning("OfficeCLI version probe failed; preview cache uses unknown renderer: %s", type(exc).__name__)
            return "unknown"

    @staticmethod
    def _payload_text(payload: dict[str, Any], *, mode: str) -> str:
        if payload.get("success") is not True or not isinstance(payload.get("data"), str):
            raise OfficePreviewMalformedError(f"OfficeCLI {mode} preview returned an invalid payload")
        return payload["data"]

    @staticmethod
    def _harden_html(raw_html: str) -> str:
        head = re.search(r"<head(?:\s[^>]*)?>", raw_html, flags=re.IGNORECASE)
        if head is None:
            raise OfficePreviewMalformedError("OfficeCLI HTML preview has no head element")
        csp = html_lib.escape(OFFICE_PREVIEW_CSP, quote=True)
        meta = f'<meta http-equiv="Content-Security-Policy" content="{csp}">'
        return raw_html[: head.end()] + meta + raw_html[head.end() :]

    @staticmethod
    def _text_fallback_html(text: str) -> str:
        csp = html_lib.escape(OFFICE_PREVIEW_CSP, quote=True)
        escaped = html_lib.escape(text)
        return (
            "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
            f'<meta http-equiv="Content-Security-Policy" content="{csp}">'
            "<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">"
            "<style>body{margin:0;padding:24px;background:#fff;color:#1f2937;font:14px/1.6 ui-monospace,monospace}"
            ".notice{margin-bottom:16px;padding:10px 12px;border:1px solid #f59e0b;border-radius:8px;background:#fffbeb}"
            "pre{white-space:pre-wrap;word-break:break-word}</style></head><body>"
            '<div class="notice">Text-only preview — layout rendering is temporarily unavailable.</div>'
            f"<pre>{escaped}</pre></body></html>"
        )

    def _checked_output_size(self, value: str) -> int:
        output_bytes = len(value.encode("utf-8"))
        if output_bytes > self.preview_max_bytes:
            raise OfficePreviewTooLargeError(
                f"Office preview output exceeds {self.preview_max_bytes} bytes"
            )
        return output_bytes

    def _load_preview_cache(
        self,
        cache_dir: Path,
        *,
        source_sha256: str,
        renderer_version: str,
    ) -> OfficePreviewResult | None:
        manifest_path = cache_dir / "manifest.json"
        html_path = cache_dir / "current.html"
        try:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            preview_html = html_path.read_text(encoding="utf-8")
        except (OSError, json.JSONDecodeError):
            return None
        if (
            manifest.get("source_sha256") != source_sha256
            or manifest.get("renderer_version") != renderer_version
            or manifest.get("preview_sha256") != sha256(preview_html.encode("utf-8")).hexdigest()
        ):
            return None
        try:
            preview_mode = str(manifest["preview_mode"])
            output_bytes = int(manifest["output_bytes"])
        except (KeyError, TypeError, ValueError):
            return None
        if preview_mode not in {"html", "text_fallback"} or output_bytes != len(preview_html.encode("utf-8")):
            return None
        if output_bytes > self.preview_max_bytes:
            return None
        return OfficePreviewResult(
            html=preview_html,
            preview_mode=preview_mode,
            source_sha256=source_sha256,
            renderer_version=renderer_version,
            cache_hit=True,
            output_bytes=output_bytes,
            fallback_reason=manifest.get("fallback_reason"),
        )

    def _save_preview_cache(self, cache_dir: Path, result: OfficePreviewResult) -> None:
        cache_dir.mkdir(parents=True, exist_ok=True)
        self._atomic_write_text(cache_dir / "current.html", result.html)
        manifest = {
            "source_sha256": result.source_sha256,
            "renderer_version": result.renderer_version,
            "preview_mode": result.preview_mode,
            "preview_sha256": sha256(result.html.encode("utf-8")).hexdigest(),
            "output_bytes": result.output_bytes,
            "fallback_reason": result.fallback_reason,
            "generated_at": self._now(),
        }
        self._atomic_write_text(
            cache_dir / "manifest.json",
            json.dumps(manifest, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        )

    @staticmethod
    def _sha256_file(path: Path) -> str:
        digest = sha256()
        with path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(chunk)
        return digest.hexdigest()

    def _require_existing_file(self, rel_path: str) -> Path:
        target = self.resolve_document_path(rel_path)
        if not target.is_file():
            raise OfficeDocumentNotFoundError(f"Document not found: {rel_path}")
        return target

    def _normalize_rel_path(self, rel_path: str) -> str:
        raw = (rel_path or "").strip()
        if not raw:
            raise OfficeDocumentPathError("Document path is required")
        path = Path(raw)
        if path.is_absolute():
            raise OfficeDocumentPathError("Absolute document paths are not allowed")
        if any(part in {"..", ""} for part in path.parts):
            raise OfficeDocumentPathError("Document path cannot escape the workspace")
        if path.parts and path.parts[0] == ".office_meta":
            raise OfficeDocumentPathError("Office sidecar metadata is not a document path")
        return path.as_posix()

    def _assert_under_workspace(self, target: Path) -> None:
        try:
            target.relative_to(self.workspace)
        except ValueError as exc:
            raise OfficeDocumentPathError("Document path cannot escape the workspace") from exc

    def _load_manifest(self, rel_path: str) -> dict[str, Any]:
        path = self.manifest_path(rel_path)
        if not path.exists():
            return {
                "path": rel_path,
                "kind": self._kind_for_path(Path(rel_path)),
                "current_version": 0,
                "created_at": self._now(),
                "updated_at": self._now(),
                "revisions": [],
                "operations": [],
            }
        try:
            manifest = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            raise OfficeDocumentError(f"Invalid office manifest JSON for {rel_path}") from exc
        if not isinstance(manifest, dict):
            raise OfficeDocumentError(f"Invalid office manifest object for {rel_path}")
        manifest.pop("active_editor_session", None)
        return manifest

    def _save_manifest(self, rel_path: str, manifest: dict[str, Any]) -> None:
        path = self.manifest_path(rel_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self._atomic_write_text(path, json.dumps(manifest, ensure_ascii=False, indent=2, sort_keys=True) + "\n")

    def _write_revision(self, rel_path: str, target: Path, *, version: int, reason: str) -> dict[str, Any]:
        meta_dir = self.manifest_path(rel_path).parent
        revisions_dir = meta_dir / "revisions"
        revisions_dir.mkdir(parents=True, exist_ok=True)
        revision_name = f"{version:06d}.bin"
        revision_path = revisions_dir / revision_name
        shutil.copyfile(target, revision_path)
        return {
            "version": version,
            "file": revision_name,
            "reason": reason,
            "created_at": self._now(),
            "size": revision_path.stat().st_size,
        }

    def _write_operation(self, rel_path: str, operations: list[dict[str, Any]]) -> str:
        meta_dir = self.manifest_path(rel_path).parent
        operations_dir = meta_dir / "operations"
        operations_dir.mkdir(parents=True, exist_ok=True)
        operation_name = f"{datetime.now(UTC).strftime('%Y%m%dT%H%M%S%f')}.json"
        operation_path = operations_dir / operation_name
        self._atomic_write_text(operation_path, json.dumps({"operations": operations}, ensure_ascii=False, indent=2))
        return str(operation_path)

    def _create_blank_document(self, target: Path, kind: str) -> None:
        with tempfile.NamedTemporaryFile(delete=False, suffix=target.suffix, dir=target.parent) as handle:
            tmp_path = Path(handle.name)
        try:
            if kind == "docx":
                from docx import Document

                document = Document()
                document.add_paragraph("")
                document.save(tmp_path)
            elif kind == "xlsx":
                from openpyxl import Workbook

                workbook = Workbook()
                workbook.save(tmp_path)
            elif kind == "pptx":
                from pptx import Presentation

                presentation = Presentation()
                presentation.slides.add_slide(presentation.slide_layouts[6])
                presentation.save(tmp_path)
            else:
                raise ValueError(f"Unsupported office document kind: {kind}")
            os.replace(tmp_path, target)
        finally:
            if tmp_path.exists():
                tmp_path.unlink()

    def _atomic_copy(self, source: Path, target: Path) -> None:
        with tempfile.NamedTemporaryFile(delete=False, suffix=target.suffix, dir=target.parent) as handle:
            tmp_path = Path(handle.name)
        try:
            shutil.copyfile(source, tmp_path)
            os.replace(tmp_path, target)
        finally:
            if tmp_path.exists():
                tmp_path.unlink()

    @staticmethod
    def _atomic_write_bytes(path: Path, content: bytes) -> None:
        with tempfile.NamedTemporaryFile(delete=False, dir=path.parent) as handle:
            tmp_path = Path(handle.name)
            handle.write(content)
        try:
            os.replace(tmp_path, path)
        finally:
            if tmp_path.exists():
                tmp_path.unlink()

    @staticmethod
    def _atomic_write_text(path: Path, content: str) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        with tempfile.NamedTemporaryFile("w", delete=False, dir=path.parent, encoding="utf-8") as handle:
            tmp_path = Path(handle.name)
            handle.write(content)
        try:
            os.replace(tmp_path, path)
        finally:
            if tmp_path.exists():
                tmp_path.unlink()

    @staticmethod
    def _kind_for_path(path: Path) -> str | None:
        return SUPPORTED_OFFICE_SUFFIXES.get(path.suffix.lower())

    @staticmethod
    def _now() -> str:
        return datetime.now(UTC).isoformat()
