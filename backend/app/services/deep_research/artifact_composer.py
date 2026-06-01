from __future__ import annotations

import html
import json
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from app.services.deep_research.plan_contract import SUPPORTED_DEEP_RESEARCH_FORMATS, normalize_contract_output_format
from app.services.office_document_service import OfficeDocumentService

SUPPORTED_DEEP_RESEARCH_OUTPUT_FORMATS: set[str] = set(SUPPORTED_DEEP_RESEARCH_FORMATS)


@dataclass(slots=True)
class ResearchDossier:
    artifact_dir: Path
    report_markdown: str
    final: dict[str, Any]
    sources: list[dict[str, Any]]
    claims: list[dict[str, Any]]
    lane_summaries: list[dict[str, Any]]
    quality_gates: dict[str, Any]
    gaps: list[str]


def compose_deep_research_artifact(workspace: Path, artifact_dir: Path, output_format: str) -> Path | None:
    normalized = normalize_contract_output_format(output_format)
    dossier = load_research_dossier(artifact_dir)
    if normalized == "markdown":
        return artifact_dir / "report.md"
    if normalized == "json":
        return artifact_dir / "final.json"
    if not dossier.report_markdown and normalized != "json":
        return None
    if normalized == "html":
        target = artifact_dir / "report.html"
        target.write_text(_compose_html(dossier), encoding="utf-8")
        return target
    if normalized == "docx":
        target = artifact_dir / "report.docx"
        _save_office_bytes(workspace, target, _compose_docx(dossier), reason="deep-research-docx-compose")
        return target
    if normalized == "xlsx":
        target = artifact_dir / "report.xlsx"
        _save_office_bytes(workspace, target, _compose_xlsx(dossier), reason="deep-research-xlsx-compose")
        return target
    if normalized == "pptx":
        target = artifact_dir / "report.pptx"
        _save_office_bytes(workspace, target, _compose_pptx(dossier), reason="deep-research-pptx-compose")
        return target
    return artifact_dir / "report.md"


def load_research_dossier(artifact_dir: Path) -> ResearchDossier:
    final = _load_json(artifact_dir / "final.json")
    return ResearchDossier(
        artifact_dir=artifact_dir,
        report_markdown=_read_text(artifact_dir / "report.md"),
        final=final,
        sources=_load_jsonl(artifact_dir / "sources.jsonl") or final.get("sources", []),
        claims=_load_jsonl(artifact_dir / "claims.jsonl") or final.get("claims", []),
        lane_summaries=_load_jsonl(artifact_dir / "lane_summaries.jsonl"),
        quality_gates=final.get("quality_gates") if isinstance(final.get("quality_gates"), dict) else {},
        gaps=final.get("gaps") if isinstance(final.get("gaps"), list) else [],
    )


def _compose_docx(dossier: ResearchDossier) -> bytes:
    from docx import Document

    document = Document()
    title = _title_from_markdown(dossier.report_markdown) or "Deep Research Report"
    document.add_heading(title, level=0)
    document.add_heading("Executive Summary", level=1)
    document.add_paragraph(str(dossier.final.get("summary") or _first_paragraph(dossier.report_markdown) or ""))
    for heading, body in _markdown_sections(dossier.report_markdown):
        if heading.casefold() in {"source ledger"}:
            continue
        document.add_heading(heading, level=1)
        for paragraph in body:
            if paragraph.startswith("- "):
                document.add_paragraph(paragraph[2:].strip(), style="List Bullet")
            else:
                document.add_paragraph(paragraph)
    document.add_heading("Evidence Appendix", level=1)
    if dossier.sources:
        table = document.add_table(rows=1, cols=5)
        headers = ["Source ID", "Title", "Publisher", "Type", "URL"]
        for index, header in enumerate(headers):
            table.rows[0].cells[index].text = header
        for source in dossier.sources[:40]:
            row = table.add_row().cells
            row[0].text = str(source.get("source_id") or "")
            row[1].text = str(source.get("title") or "")
            row[2].text = str(source.get("publisher") or "")
            row[3].text = str(source.get("source_type") or "")
            row[4].text = str(source.get("url") or "")
    if dossier.gaps:
        document.add_heading("Gaps", level=2)
        for gap in dossier.gaps:
            document.add_paragraph(str(gap), style="List Bullet")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _compose_xlsx(dossier: ResearchDossier) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.append(["Field", "Value"])
    summary.append(["status", dossier.final.get("status")])
    summary.append(["summary", dossier.final.get("summary")])
    summary.append(["source_count", len(dossier.sources)])
    summary.append(["claim_count", len(dossier.claims)])

    sources = workbook.create_sheet("Sources")
    sources.append(["source_id", "title", "publisher", "url", "source_type", "evidence_tier", "evidence_grade"])
    for source in dossier.sources:
        sources.append(
            [
                source.get("source_id"),
                source.get("title"),
                source.get("publisher"),
                source.get("url"),
                source.get("source_type"),
                source.get("evidence_tier"),
                source.get("evidence_grade"),
            ]
        )

    claims = workbook.create_sheet("Claims")
    claims.append(["claim_id", "text", "status", "source_ids", "evidence", "notes"])
    for claim in dossier.claims:
        claims.append(
            [
                claim.get("claim_id"),
                claim.get("text"),
                claim.get("status"),
                ", ".join(str(item) for item in claim.get("source_ids", []) if item),
                claim.get("evidence"),
                claim.get("notes"),
            ]
        )

    lanes = workbook.create_sheet("Lane Coverage")
    lanes.append(["lane_id", "label", "source_count", "evidence_strength", "key_findings", "limitations"])
    for lane in dossier.lane_summaries:
        lanes.append(
            [
                lane.get("lane_id"),
                lane.get("label"),
                lane.get("source_count"),
                lane.get("evidence_strength"),
                "\n".join(str(item) for item in lane.get("key_findings", []) if item),
                "\n".join(str(item) for item in lane.get("limitations", []) if item),
            ]
        )

    gates = workbook.create_sheet("Quality Gates")
    gates.append(["gate", "status"])
    for gate, status in dossier.quality_gates.items():
        gates.append([gate, status])
    gaps = workbook.create_sheet("Risks And Gaps")
    gaps.append(["gap"])
    for gap in dossier.gaps:
        gaps.append([str(gap)])

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _compose_pptx(dossier: ResearchDossier) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    title = _title_from_markdown(dossier.report_markdown) or "Deep Research"
    sections = dict(_markdown_sections(dossier.report_markdown))
    _add_slide(presentation, "Decision Question", [title, str(dossier.final.get("summary") or "")])
    _add_slide(presentation, "Executive Thesis", _body_lines(sections.get("Executive Thesis", []), limit=5))
    evidence_lines = []
    for claim in dossier.claims[:5]:
        evidence_lines.append(str(claim.get("text") or claim.get("evidence") or ""))
    if not evidence_lines:
        evidence_lines = [f"{len(dossier.sources)} sources reviewed", f"{len(dossier.claims)} claims extracted"]
    _add_slide(presentation, "Evidence Matrix", evidence_lines[:6])
    _add_slide(presentation, "Risks And Gaps", [str(gap) for gap in dossier.gaps[:6]] or ["No major gap recorded."])
    implication_lines = _body_lines(sections.get("Strategic Implications", []), limit=5)
    if implication_lines:
        _add_slide(presentation, "Strategic Implications", implication_lines)
    source_lines = [
        f"{source.get('source_id')}: {source.get('title') or source.get('publisher')}"
        for source in dossier.sources[:8]
    ]
    _add_slide(presentation, "Source Appendix", source_lines or ["No sources recorded."])
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _compose_html(dossier: ResearchDossier) -> str:
    sections = _markdown_sections(dossier.report_markdown)
    body = [f"<h1>{html.escape(_title_from_markdown(dossier.report_markdown) or 'Deep Research Report')}</h1>"]
    body.append(f"<p><strong>Summary:</strong> {html.escape(str(dossier.final.get('summary') or ''))}</p>")
    for heading, paragraphs in sections:
        body.append(f"<section><h2>{html.escape(heading)}</h2>")
        for paragraph in paragraphs:
            body.append(f"<p>{html.escape(paragraph)}</p>")
        body.append("</section>")
    if dossier.sources:
        body.append("<section><h2>Evidence Cards</h2>")
        for source in dossier.sources:
            title = html.escape(str(source.get("title") or source.get("url") or "Source"))
            url = html.escape(str(source.get("url") or ""))
            body.append(f"<article><strong>{title}</strong><br><a href=\"{url}\">{url}</a></article>")
        body.append("</section>")
    return (
        "<!doctype html><html><head><meta charset=\"utf-8\"><title>Deep Research Report</title>"
        "<style>body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;line-height:1.55;"
        "max-width:960px;margin:40px auto;padding:0 24px;color:#111827}"
        "section{margin:28px 0}article{border-top:1px solid #ddd;padding:10px 0}</style></head><body>"
        + "\n".join(body)
        + "</body></html>\n"
    )


def _add_slide(presentation: Any, title: str, body_lines: list[str]) -> None:
    layout = presentation.slide_layouts[0] if not presentation.slides else presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = title[:120]
    placeholders = [shape for shape in slide.placeholders if shape.placeholder_format.idx != 0]
    if placeholders:
        placeholders[0].text = "\n".join(line for line in body_lines if line)[:1200]


def _markdown_sections(markdown: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_heading = ""
    current_body: list[str] = []
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("# "):
            continue
        if stripped.startswith("## "):
            if current_heading:
                sections.append((current_heading, current_body))
            current_heading = stripped[3:].strip()
            current_body = []
            continue
        if stripped and current_heading:
            current_body.append(stripped)
    if current_heading:
        sections.append((current_heading, current_body))
    return sections


def _body_lines(lines: list[str], *, limit: int) -> list[str]:
    return [line[2:].strip() if line.startswith("- ") else line for line in lines if line.strip()][:limit]


def _title_from_markdown(markdown: str) -> str:
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("# "):
            return stripped[2:].strip()
    return ""


def _first_paragraph(markdown: str) -> str:
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped and not stripped.startswith("#"):
            return stripped
    return ""


def _save_office_bytes(workspace: Path, target: Path, content: bytes, *, reason: str) -> None:
    rel_path = _relative(workspace, target)
    OfficeDocumentService(workspace).atomic_save_bytes(
        rel_path,
        content,
        reason=reason,
        require_no_active_editor=True,
    )


def _relative(workspace: Path, path: Path) -> str:
    try:
        return path.resolve().relative_to(workspace.resolve()).as_posix()
    except ValueError:
        return path.as_posix()


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except OSError:
        return ""


def _load_json(path: Path) -> dict[str, Any]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return payload if isinstance(payload, dict) else {}


def _load_jsonl(path: Path) -> list[dict[str, Any]]:
    try:
        lines = path.read_text(encoding="utf-8").splitlines()
    except OSError:
        return []
    payloads: list[dict[str, Any]] = []
    for line in lines:
        if not line.strip():
            continue
        try:
            payload = json.loads(line)
        except json.JSONDecodeError:
            continue
        if isinstance(payload, dict):
            payloads.append(payload)
    return payloads
