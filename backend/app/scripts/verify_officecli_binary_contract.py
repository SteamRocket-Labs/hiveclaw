from __future__ import annotations

import argparse
import html as html_lib
import json
import shutil
import tempfile
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.services.office_document_service import (
    OFFICE_PREVIEW_CSP,
    OfficeDocumentService,
    OfficePreviewMalformedError,
    extract_officecli_text_payload,
)
from app.services.officecli_adapter import OfficeCLIAdapter


def _write_contract_fixtures(root: Path) -> dict[str, Path]:
    docx_path = root / "contract.docx"
    document = Document()
    document.add_heading("OfficeCLI contract", level=1)
    document.add_paragraph("DOCX preview evidence")
    document.add_table(rows=1, cols=2).rows[0].cells[0].text = "table"
    document.save(docx_path)

    xlsx_path = root / "contract.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet["A1"] = "OfficeCLI contract"
    sheet["B1"] = "=1+1"
    workbook.save(xlsx_path)

    pptx_path = root / "contract.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    textbox.text_frame.text = "OfficeCLI contract PPTX preview evidence"
    presentation.save(pptx_path)

    return {"docx": docx_path, "xlsx": xlsx_path, "pptx": pptx_path}


def _valid_html_view_payload(payload: dict[str, Any]) -> bool:
    return payload.get("success") is True and isinstance(payload.get("data"), str) and bool(payload["data"].strip())


def _valid_text_view_payload(payload: dict[str, Any], *, office_format: str) -> bool:
    try:
        return bool(extract_officecli_text_payload(payload, office_format=office_format).strip())
    except OfficePreviewMalformedError:
        return False


def _preview_contains_csp(rendered_html: str) -> bool:
    escaped_csp = html_lib.escape(OFFICE_PREVIEW_CSP, quote=True)
    return (
        'http-equiv="Content-Security-Policy"' in rendered_html
        and f'content="{escaped_csp}"' in rendered_html
    )


def _contract_format_ok(result: dict[str, Any]) -> bool:
    return (
        result.get("html") is True
        and result.get("text") is True
        and result.get("service_preview_mode") == "html"
        and result.get("csp") is True
    )


def verify_officecli_binary_contract(*, binary: str | Path) -> dict[str, Any]:
    """Exercise the deployed OfficeCLI binary through the exact Hive preview contract."""

    adapter = OfficeCLIAdapter(binary=str(binary), binary_sha256="", timeout_seconds=60)
    with tempfile.TemporaryDirectory(prefix="hive-officecli-contract-") as temp_dir:
        root = Path(temp_dir)
        fixtures = _write_contract_fixtures(root)
        service = OfficeDocumentService(root, adapter=adapter, preview_max_bytes=25 * 1024 * 1024)
        format_results: dict[str, dict[str, Any]] = {}
        for office_format, path in fixtures.items():
            html_payload = adapter.run_view(path, mode="html", cwd=root)
            text_payload = adapter.run_view(path, mode="text", cwd=root)
            preview = service.render_preview(path.name)
            format_results[office_format] = {
                "html": _valid_html_view_payload(html_payload),
                "text": _valid_text_view_payload(text_payload, office_format=office_format),
                "service_preview_mode": preview.preview_mode,
                "csp": _preview_contains_csp(preview.html),
                "output_bytes": preview.output_bytes,
            }

    failed_formats = sorted(
        office_format for office_format, result in format_results.items() if not _contract_format_ok(result)
    )
    if failed_formats:
        raise RuntimeError(f"OfficeCLI binary contract failed for: {', '.join(failed_formats)}")

    return {
        "status": "ok",
        "version": adapter.version(),
        "formats": format_results,
    }


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Verify the deployed OfficeCLI HTML/text preview contract")
    parser.add_argument("--binary", default=shutil.which("officecli") or "officecli")
    return parser.parse_args()


def main() -> int:
    args = _parse_args()
    try:
        report = verify_officecli_binary_contract(binary=args.binary)
    except Exception as exc:
        print(json.dumps({"status": "error", "error_type": type(exc).__name__}, sort_keys=True))
        return 1
    print(json.dumps(report, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
